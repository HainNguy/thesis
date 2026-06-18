from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x33, 0x66)   # titles, diagram boxes
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)   # subtle rule / accent
MILESTONE_GOLD = RGBColor(0xC8, 0x96, 0x20)

FONT_TITLE  = "Calibri"
FONT_BODY   = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def txbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))

def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(14)
            run.font.name = FONT_BODY

def title_box(slide, text, top=0.35):
    tb = txbox(slide, 0.5, top, 12.3, 0.85)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(36)
    run.font.bold  = True
    run.font.name  = FONT_TITLE
    run.font.color.rgb = DARK_BLUE
    # thin rule below title
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(top + 0.88),
        Inches(12.3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = DARK_BLUE
    line.line.fill.background()
    return tb

def bullet_box(slide, bullets, left=0.5, top=1.55, width=12.3, height=5.4,
               size=24, bold_first=False):
    tb = txbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.name  = FONT_BODY
        run.font.color.rgb = DARK_TEXT
        run.font.bold  = (bold_first and i == 0)

def add_rect(slide, left, top, width, height,
             fill=DARK_BLUE, text="", font_size=16):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(font_size)
        run.font.bold  = True
        run.font.name  = FONT_BODY
        run.font.color.rgb = WHITE
    return shape

def add_arrow(slide, left, top, width=0.4, height=0.03):
    """Horizontal arrow as a thin dark rectangle + triangle approximation."""
    # shaft
    shaft = slide.shapes.add_shape(
        1, Inches(left), Inches(top + 0.22),
        Inches(width - 0.15), Inches(0.06)
    )
    shaft.fill.solid()
    shaft.fill.fore_color.rgb = DARK_TEXT
    shaft.line.fill.background()
    # arrowhead (right-pointing triangle via very short wide rect)
    head = slide.shapes.add_shape(
        5,  # right triangle — MSO_SHAPE_TYPE is not clean, use isoceles triangle
        Inches(left + width - 0.18), Inches(top + 0.12),
        Inches(0.18), Inches(0.26)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = DARK_TEXT
    head.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 · Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()

# Full-width dark header band
band = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(3.2))
band.fill.solid()
band.fill.fore_color.rgb = DARK_BLUE
band.line.fill.background()

tb = txbox(s1, 0.6, 0.55, 12.1, 2.0)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Benchmarking Tabular Data Augmentation\nTechniques for Deep Clustering"
run.font.size  = Pt(36)
run.font.bold  = True
run.font.name  = FONT_TITLE
run.font.color.rgb = WHITE

# Sub-info
sub = txbox(s1, 0.6, 3.5, 12.1, 2.0)
tf2 = sub.text_frame
tf2.word_wrap = True
for line, sz, bold in [
    ("Hai Nguyen  ·  B.Sc. Computer Science  ·  LMU Munich", 22, False),
    ("Supervised by Mamdouh Aljoud  ·  DBS Group", 20, False),
    ("Preliminary Presentation  ·  12 June 2026", 20, False),
]:
    p2 = tf2.paragraphs[0] if line == "Hai Nguyen  ·  B.Sc. Computer Science  ·  LMU Munich" \
         else tf2.add_paragraph()
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = text = line
    run2.font.size  = Pt(sz)
    run2.font.bold  = bold
    run2.font.name  = FONT_BODY
    run2.font.color.rgb = DARK_TEXT

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 · Introduction
# ══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
title_box(s2, "Introduction")
bullet_box(s2, [
    "Clustering is useful when labels are unavailable — many real-world datasets have no ground-truth class information.",
    "Classical methods struggle on tabular data: mixed feature types, noise, and nonlinear relationships make raw features poor inputs for k-means.",
    "Deep clustering addresses this by learning a latent representation that is easier to cluster — a neural network maps raw features to a compact space first.",
])
set_notes(s2,
    "Classical k-means on raw features fails when features are on different scales or when clusters are not "
    "linearly separable. Deep clustering uses a neural network to first learn a better space."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 · Motivation
# ══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
title_box(s3, "Motivation")
bullet_box(s3, [
    "Contrastive learning builds robust representations by pulling augmented views of the same sample together and pushing different samples apart.",
    "In computer vision, standard augmentations are well-established: random crop, flip, colour jitter — each preserves semantic content.",
    "For tabular data, no consensus exists: some augmentations may destroy the cluster structure rather than preserve it, making the choice non-trivial.",
])
set_notes(s3,
    "The core problem is that unlike images, there is no obvious 'rotate a tabular row' operation. "
    "The augmentation has to be designed carefully so it creates a different view of the same sample "
    "without changing which cluster it belongs to."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 · Research Question
# ══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
title_box(s4, "Research Question")

# Centred large-font question in a shaded box
box = s4.shapes.add_shape(1, Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.0))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_GREY
box.line.color.rgb = DARK_BLUE

tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)
run = p.add_run()
run.text = (
    "Which tabular augmentation techniques are most effective for learning "
    "cluster-friendly representations with contrastive learning, and how do "
    "they impact deep clustering performance?"
)
run.font.size  = Pt(26)
run.font.bold  = True
run.font.name  = FONT_TITLE
run.font.color.rgb = DARK_BLUE

set_notes(s4,
    "This is the exact research question from the approved proposal. The key phrase is cluster-friendly — "
    "we are not just measuring representation quality, we care whether the clusters are preserved."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 · Approach  (flow diagram + technique bullets)
# ══════════════════════════════════════════════════════════════════════════════
s5 = add_slide()
title_box(s5, "Approach")

# 5-box flow diagram  ─────────────────────────────────────────────────────────
boxes = [
    "Raw\nTabular\nData",
    "Augmentation\n(masking /\nnoise /\nswapping)",
    "MLP\nEncoder",
    "k-means\nClustering",
    "Evaluate\nNMI / ARI",
]
box_w, box_h = 1.9, 1.15
gap         = 0.32
start_left  = 0.55
top_row     = 1.55

for i, label in enumerate(boxes):
    bx = start_left + i * (box_w + gap)
    add_rect(s5, bx, top_row, box_w, box_h, fill=DARK_BLUE, text=label, font_size=15)
    if i < len(boxes) - 1:
        add_arrow(s5, bx + box_w + 0.01, top_row + 0.3, width=gap + 0.12)

# Augmentation technique bullets (below diagram)
bullet_box(s5, [
    "Feature masking (SCARF-style)  —  randomly corrupt a fraction of features by replacing them with values drawn from the marginal distribution",
    "Gaussian noise  —  add small random noise to numerical features; leaves categorical features unchanged",
    "Feature swapping  —  swap feature values between randomly selected samples in the same batch",
], left=0.5, top=2.95, width=12.3, height=3.8, size=20)

set_notes(s5,
    "The pipeline is the same regardless of which augmentation we use — that is the point of a benchmark. "
    "Only the augmentation block changes; everything else is fixed."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 · Timeline
# ══════════════════════════════════════════════════════════════════════════════
s6 = add_slide()
title_box(s6, "Timeline")

phases = [
    ("Preparation",        "until 14 Jun 2026",   "Tool onboarding: Python, PyTorch"),
    ("Phase A  Literature","15–28 Jun  (Wk 1–2)", "Literature review & methodology design"),
    ("Phase B  Implementation","29 Jun–19 Jul  (Wk 3–5)", "Build benchmark pipeline"),
    ("Phase C  Experiments","20 Jul–2 Aug  (Wk 6–7)", "Run experiments & sensitivity studies"),
    ("Phase D–E  Writing & Submission","3–23 Aug  (Wk 8–10)", "Write thesis · submit 23 Aug 2026"),
]
milestones = {1: "M1 · PyTorch ready (21 Jun)", 3: "M3 · MVP pipeline (12 Jul)", 4: "M6 · Submission (23 Aug)"}

row_h   = 0.72
col_w   = [3.2, 2.8, 4.4, 1.6]
left0   = 0.5
top0    = 1.5

# Header row
for ci, (text, w) in enumerate(zip(["Phase", "Period", "Focus", "Milestone"], col_w)):
    cl = left0 + sum(col_w[:ci]) + ci * 0.05
    hdr = s6.shapes.add_shape(1, Inches(cl), Inches(top0), Inches(w), Inches(0.46))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()
    tf = hdr.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.name = FONT_BODY
    run.font.color.rgb = WHITE

for ri, (phase, period, focus) in enumerate(phases):
    row_top = top0 + 0.46 + ri * (row_h + 0.04)
    bg = LIGHT_GREY if ri % 2 == 0 else WHITE
    milestone_text = milestones.get(ri, "")
    for ci, (text, w) in enumerate(zip([phase, period, focus, milestone_text], col_w)):
        cl = left0 + sum(col_w[:ci]) + ci * 0.05
        cell = s6.shapes.add_shape(1, Inches(cl), Inches(row_top), Inches(w), Inches(row_h))
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.line.color.rgb = LIGHT_GREY
        tf = cell.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci < 3 else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(15)
        run.font.bold  = (ci == 0)
        run.font.name  = FONT_BODY
        run.font.color.rgb = MILESTONE_GOLD if (ci == 3 and text) else DARK_TEXT

set_notes(s6,
    "The official thesis period is only 10 weeks, so all tool learning — including PyTorch — is deliberately "
    "moved into the preparation phase before June 15 to keep the full official period for actual research."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 · Expected Challenges
# ══════════════════════════════════════════════════════════════════════════════
s7 = add_slide()
title_box(s7, "Expected Challenges")
bullet_box(s7, [
    "PyTorch learning curve — contrastive loss (NT-Xent/InfoNCE) and training loop are new; fully addressed in the preparation phase before the official start date.",
    "Mixed feature types — Gaussian noise suits numerical features; categorical features require masking or swapping instead; augmentation must handle both correctly.",
    "Dataset selection — NMI and ARI require ground-truth labels; selecting 3–4 appropriate OpenML/UCI datasets with labels is a deliberate design choice.",
    "Negative transfer risk — augmentations that help representation learning in vision may destroy tabular cluster boundaries; identifying failure cases is a core research goal.",
], size=22)
set_notes(s7,
    "I am naming these challenges explicitly because they are already reflected in the timeline design — "
    "the preparation phase and the sensitivity study week exist specifically because of challenges 1 and 4."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 · Summary
# ══════════════════════════════════════════════════════════════════════════════
s8 = add_slide()
title_box(s8, "Summary")
bullet_box(s8, [
    "Research question: which tabular augmentation techniques best support cluster-friendly representations with contrastive learning?",
    "Approach: benchmark three augmentations (feature masking, Gaussian noise, feature swapping) through a fixed MLP → k-means pipeline, evaluated by NMI and ARI.",
    "Timeline: preparation complete by 14 Jun · official phase 15 Jun–23 Aug 2026 · six verifiable milestones.",
    "Thank you — happy to take questions.",
], size=23)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/bernstein/Library/Mobile Documents/com~apple~CloudDocs/Studium/Bachelor Arbeit/Preliminary_Presentation_Hai_Nguyen.pptx"
prs.save(out)
print(f"Saved → {out}")
